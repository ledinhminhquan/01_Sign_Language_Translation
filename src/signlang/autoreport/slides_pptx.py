"""Generate the submission slides.pptx (python-pptx) - ~12 concise slides for signlang.
Degrades to a Markdown outline if python-pptx is unavailable.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from ..config import AppConfig, run_dir
from ..logging_utils import get_logger
from . import charts as charts_mod
from .artifact_loader import (agent_recog, agent_trans, baseline_recog, load_artifacts, seg_f1)

logger = get_logger(__name__)


def _pct(v, scale=100.0):
    return f"{v*scale:.1f}%" if isinstance(v, (int, float)) and not isinstance(v, bool) else "?"


def _num(v):
    return f"{v:.1f}" if isinstance(v, (int, float)) and not isinstance(v, bool) else "?"


def _slides(cfg: AppConfig, arts: Dict[str, Any]) -> List[Tuple[str, List[str]]]:
    acc = agent_recog(arts, "gloss_accuracy")
    mf = baseline_recog(arts, "most_frequent_gloss", "gloss_accuracy")
    bleu = agent_trans(arts, "bleu")
    res = (f"gloss accuracy {_pct(acc)} vs most-frequent {_pct(mf)}; BLEU {_num(bleu)}"
           if acc is not None else "train + evaluate to populate results")
    return [
        ("Sign Language Translation",
         [f"{cfg.author} - Student {cfg.student_id}", "NLP in Industry - Final Assignment",
          "Sign-language video / pose -> spoken text (via gloss)",
          "Frozen pose front-end + a TRAINABLE seq2seq core",
          "Agent segments, recognizes, verifies, and abstains"]),
        ("Business Problem & Motivation",
         ["Accessibility: bridge Deaf signers and non-signers",
          "Assistive, NOT a replacement for human interpreters",
          "Video/pose -> gloss recognition -> text translation (Sign2Gloss2Text)",
          "Only the seq2seq core is trained; the pose front-end is algorithmic"]),
        ("Proposed Solution",
         ["MediaPipe Holistic pose (frozen) / SeedPoseEngine offline",
          "Motion-based sign segmentation",
          "Per-segment gloss recognizer + gloss->text translator",
          "Confidence gating + verification + abstention"]),
        ("System Architecture",
         ["video -> pose keypoints -> segment signs",
          "recognize gloss (D3) -> translate to text (D4)",
          "abstain on unclear / out-of-vocabulary signing (D5)",
          "Runs fully offline (SeedPose + numpy centroid + lexicon) for tests/CI"]),
        ("Data (Synthetic Pose + Real Corpora)",
         ["NO permissive continuous-SLT corpus exists on HF -> synthetic generator is PRIMARY",
          "Synthetic: deterministic keypoint trajectories + embedded gold gloss/text",
          "Real (flagged): iSign NC+gated, How2Sign NC, WLASL 'other'",
          "Permissive smoke test: Sigurdur/icelandic-sign-language (Apache)"]),
        ("The Pose Front-End + Seq2Seq Core",
         ["Front-end: MediaPipe Holistic (algorithmic) / xclip-base (MIT) - frozen",
          "Recognizer: transformer over pose segments (numpy centroid offline)",
          "Translator: t5-small (Apache) / m2m100_418M (MIT, reused P13/P14)",
          "No pretrained SLT checkpoint loads cleanly -> train our own small core"]),
        ("Metrics",
         ["Recognition: gloss WER + accuracy + sequence exact-match",
          "Translation: BLEU-1..4 (BLEU-4 headline) + chrF + WER",
          "Segmentation: boundary-F1; abstain rate",
          "Caveat: automatic SLT metrics are unreliable (Yazdani 2025) - flagged"]),
        ("The 5-Decision Agent",
         ["D1 ingest gate - D2 motion segmentation",
          "D3 per-sign confidence gate",
          "D4 translate + round-trip verify",
          "D5 ABSTAIN on out-of-vocabulary / unclear signing"]),
        ("Evaluation Results",
         [res,
          f"segmentation boundary-F1 {_pct(seg_f1(arts))}",
          "vs most-frequent / random / identity-translate baselines",
          "Confidence gating + verification + abstention = the value-add"]),
        ("Deployment Overview",
         ["FastAPI POST /translate (seed|frames -> glosses + text + confidence + abstain)",
          "Gradio demo (pick a signed sentence -> see glosses + text)",
          "Docker (mediapipe + ffmpeg + libGL); offline SeedPose + centroid on CPU",
          "Metadata-only job logging (pose video is biometric)"]),
        ("Continual Learning, Monitoring & Ethics",
         ["Monitor: abstain/low-confidence rate + recognition-confidence drift",
          "New signs/signers -> expand vocabulary + re-train the recognizer",
          "Privacy: sign video = biometric + identifying -> consent, edge, no retention",
          "Representation bias acute -> abstain + show confidence + human in the loop"]),
        ("Key Takeaways & Future Work",
         ["A segmenting, confidence-gating, abstaining sign-translation pipeline",
          "Trainable seq2seq core beats most-frequent / random floors",
          "Future: neural CSLR (CTC), real corpora (with consent), more sign languages",
          "Future: gloss-free Sign2Text, signer-independence, Deaf-community co-design"]),
    ]


def generate_slides(cfg: AppConfig, title: Optional[str] = None, author: Optional[str] = None,
                    out_path: Optional[str] = None) -> str:
    arts = load_artifacts(cfg)
    out_path = Path(out_path) if out_path else run_dir() / "report" / "slides.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    slides = _slides(cfg, arts)
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except Exception as exc:
        logger.warning("python-pptx unavailable (%s); writing markdown outline", exc)
        md = "\n\n".join(f"## {t}\n" + "\n".join(f"- {b}" for b in bs) for t, bs in slides)
        alt = out_path.with_suffix(".md")
        alt.write_text(md, encoding="utf-8")
        return str(alt)

    try:
        chart = charts_mod.recognition_chart(arts, run_dir() / "report" / "slide_recog.png")
    except Exception:
        chart = None
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    accent = RGBColor(0x2B, 0x6C, 0xB0)
    for i, (t, bullets) in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
        tf = bar.text_frame; tf.text = t
        tf.paragraphs[0].font.size = Pt(28); tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.5),
                                        Inches(8.3 if (i == 8 and chart) else 12), Inches(5.4))
        bt = body.text_frame; bt.word_wrap = True
        for j, bp in enumerate(bullets):
            p = bt.paragraphs[0] if j == 0 else bt.add_paragraph()
            p.text = "-  " + bp; p.font.size = Pt(20); p.space_after = Pt(10)
        if i == 8 and chart:
            slide.shapes.add_picture(str(chart), Inches(8.9), Inches(1.7), width=Inches(4.0))
        foot = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4))
        foot.text_frame.text = f"{title or cfg.project_title} - {author or cfg.author} ({cfg.student_id})"
        foot.text_frame.paragraphs[0].font.size = Pt(9)
    prs.save(str(out_path))
    logger.info("Slides -> %s", out_path)
    return str(out_path)


__all__ = ["generate_slides"]
